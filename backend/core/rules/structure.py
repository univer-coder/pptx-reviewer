from .base import BaseRule, Issue, Style

SLIDE_COUNT_WARNING = 20


def _get_title_shape(slide):
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                return shape
    except Exception:
        pass
    return None


class MissingTitleRule(BaseRule):
    rule_id = "missing_title"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        shape = _get_title_shape(slide)
        if shape is None or not shape.text_frame.text.strip():
            return [Issue(
                rule_id=self.rule_id,
                severity="error",
                auto_fixable=False,
                message="タイトルが設定されていません",
                suggestion="アクションタイトル（結論を示す文章）を入力してください",
                details={},
            )]
        return []


class TitleLinesRule(BaseRule):
    rule_id = "title_lines"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        shape = _get_title_shape(slide)
        if shape is None or not shape.text_frame.text.strip():
            return []
        lines = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
        if lines > 2:
            return [Issue(
                rule_id=self.rule_id,
                severity="warning",
                auto_fixable=False,
                message=f"タイトルが{lines}行あります",
                suggestion="タイトルは1行（最大2行）に収めてください",
                details={"lines": lines},
            )]
        return []


class SlideCountRule(BaseRule):
    rule_id = "slide_count"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        count = len(presentation.slides)
        if count > SLIDE_COUNT_WARNING:
            return [Issue(
                rule_id=self.rule_id,
                severity="info",
                auto_fixable=False,
                message=f"スライド枚数が多い可能性があります（{count}枚）",
                suggestion="1分〜1.5分／枚を基準に枚数を見直してください",
                details={"count": count, "threshold": SLIDE_COUNT_WARNING},
            )]
        return []


class Chart3DRule(BaseRule):
    rule_id = "chart_3d"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        issues = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.CHART:
                continue
            try:
                chart = shape.chart
                if chart.chart_type is not None:
                    ct_val = int(chart.chart_type)
                    # 3D chart type values: 3DBar=60, 3DLine=61, 3DPie=70など
                    THREE_D_TYPES = {
                        -4100, 4, 5, 6, 14, 15, 16, 17, 18, 19,
                        56, 57, 58, 59, 60, 61, 62, 63, 64, 65,
                        70, 71, 72, 73, 74, 75, 76,
                    }
                    if ct_val in THREE_D_TYPES:
                        issues.append(Issue(
                            rule_id=self.rule_id,
                            severity="warning",
                            auto_fixable=False,
                            message="3Dグラフが使用されています",
                            suggestion="3D効果を排除し、2Dグラフに変更してください",
                            details={"chart_type": ct_val},
                        ))
            except Exception:
                continue
        return issues


class PieChartRule(BaseRule):
    rule_id = "pie_chart"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.chart import XL_CHART_TYPE
        issues = []
        PIE_TYPES = {5, 6, 66, 67, 68, 69, 70, 71}
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.CHART:
                continue
            try:
                ct_val = int(shape.chart.chart_type)
                if ct_val in PIE_TYPES:
                    issues.append(Issue(
                        rule_id=self.rule_id,
                        severity="info",
                        auto_fixable=False,
                        message="円グラフが使用されています",
                        suggestion="比較には横棒グラフや帯グラフの使用を推奨します",
                        details={},
                    ))
            except Exception:
                continue
        return issues
