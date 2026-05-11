from .base import BaseRule, Issue, Style

TEXT_LENGTH_MAX = 200
LINE_COUNT_MAX = 8


def _is_title(shape) -> bool:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        return shape.is_placeholder and shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        )
    except Exception:
        return False


def _has_text(shape) -> bool:
    return shape.has_text_frame and shape.text_frame.text.strip()


class TextLengthRule(BaseRule):
    rule_id = "text_length"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        if style == Style.VISUAL:
            return []
        total = sum(
            len(shape.text_frame.text)
            for shape in slide.shapes
            if _has_text(shape)
        )
        if total > TEXT_LENGTH_MAX:
            return [Issue(
                rule_id=self.rule_id,
                severity="warning",
                auto_fixable=False,
                message=f"スライドの文字数が多すぎます（現在: {total}文字）",
                suggestion=f"{TEXT_LENGTH_MAX}文字以内に収めてください",
                details={"current": total, "maximum": TEXT_LENGTH_MAX},
            )]
        return []


class LineCountRule(BaseRule):
    rule_id = "line_count"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        if style == Style.VISUAL:
            return []
        issues = []
        for shape in slide.shapes:
            if _is_title(shape) or not _has_text(shape):
                continue
            line_count = len(shape.text_frame.paragraphs)
            if line_count > LINE_COUNT_MAX:
                issues.append(Issue(
                    rule_id=self.rule_id,
                    severity="warning",
                    auto_fixable=False,
                    message=f"テキストボックスの行数が多すぎます（現在: {line_count}行）",
                    suggestion=f"{LINE_COUNT_MAX}行以内に収めてください",
                    details={"current": line_count, "maximum": LINE_COUNT_MAX},
                ))
        return issues
