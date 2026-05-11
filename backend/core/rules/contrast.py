import math
from typing import Optional, Tuple
from pptx.dml.color import RGBColor
from .base import BaseRule, Issue, Style

WCAG_AA_RATIO = 4.5


def _relative_luminance(r: int, g: int, b: int) -> float:
    def channel(c: int) -> float:
        s = c / 255
        return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def _contrast_ratio(c1: tuple, c2: tuple) -> float:
    l1 = _relative_luminance(*c1)
    l2 = _relative_luminance(*c2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _rgb_from_shape_fill(shape) -> Optional[Tuple]:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        from pptx.enum.dml import MSO_THEME_COLOR
        fg = fill.fore_color
        if fg.rgb:
            c = fg.rgb
            return (c.red, c.green, c.blue)
    except Exception:
        return None


class ColorContrastRule(BaseRule):
    rule_id = "color_contrast"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        issues = []
        # スライド背景色を取得（デフォルト白）
        bg_color = (255, 255, 255)
        try:
            bg = slide.background.fill
            if bg.fore_color and bg.fore_color.rgb:
                c = bg.fore_color.rgb
                bg_color = (c.red, c.green, c.blue)
        except Exception:
            pass

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    try:
                        fc = run.font.color
                        if fc.type is None:
                            continue
                        rgb = fc.rgb
                        text_color = (rgb.red, rgb.green, rgb.blue)
                    except Exception:
                        continue

                    ratio = _contrast_ratio(text_color, bg_color)
                    if ratio < WCAG_AA_RATIO:
                        issues.append(Issue(
                            rule_id=self.rule_id,
                            severity="error",
                            auto_fixable=False,
                            message=f"文字色と背景色のコントラスト比が不足しています（比率: {ratio:.2f}:1）",
                            suggestion=f"WCAG AA 基準の {WCAG_AA_RATIO}:1 以上を確保してください",
                            details={"ratio": round(ratio, 2), "minimum": WCAG_AA_RATIO},
                        ))
        return issues


class ThreeColorRule(BaseRule):
    rule_id = "three_color_rule"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        base_colors = {
            (255, 255, 255), (242, 242, 242), (0, 0, 0),
            (31, 31, 31), (26, 26, 26),
        }
        accent_colors: set[tuple] = set()

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        fc = run.font.color
                        if fc.type is None:
                            continue
                        rgb = fc.rgb
                        color = (rgb.red, rgb.green, rgb.blue)
                        if color not in base_colors:
                            accent_colors.add(color)
                    except Exception:
                        continue

        if len(accent_colors) > 1:
            return [Issue(
                rule_id=self.rule_id,
                severity="warning",
                auto_fixable=False,
                message=f"アクセントカラーが複数使われています（{len(accent_colors)}色）",
                suggestion="アクセントカラーは1色に統一してください",
                details={"accent_count": len(accent_colors)},
            )]
        return []
