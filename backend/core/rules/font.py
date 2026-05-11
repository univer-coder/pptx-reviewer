from pptx.util import Pt
from .base import BaseRule, BaseFixableRule, Issue, Style

RECOMMENDED_PREFIXES = ("Meiryo", "メイリオ", "BIZ UDP")
FALLBACK_FONT = "Meiryo"
TITLE_MIN_PT = 24
TITLE_MAX_PT = 28
BODY_MIN_PT = 14


def _is_title(shape) -> bool:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        return shape.is_placeholder and shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        )
    except Exception:
        return False


def _has_text(shape) -> bool:
    return shape.has_text_frame and shape.text_frame.text.strip()


def _all_run_sizes(shape):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                yield run.font.size.pt


class TitleFontSizeRule(BaseFixableRule):
    rule_id = "title_font_size"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        for shape in slide.shapes:
            if not _is_title(shape) or not _has_text(shape):
                continue
            sizes = [s for s in _all_run_sizes(shape) if s < TITLE_MIN_PT]
            if sizes:
                pt = round(min(sizes), 1)
                return [Issue(
                    rule_id=self.rule_id,
                    severity="warning",
                    auto_fixable=True,
                    message=f"タイトルのフォントサイズが小さすぎます（最小: {pt}pt）",
                    suggestion=f"{TITLE_MIN_PT}pt 以上を推奨します",
                    details={"current": pt, "minimum": TITLE_MIN_PT},
                )]
        return []

    def fix(self, slide, presentation, style: Style) -> None:
        for shape in slide.shapes:
            if not _is_title(shape) or not _has_text(shape):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None and run.font.size < Pt(TITLE_MIN_PT):
                        run.font.size = Pt(TITLE_MIN_PT)


class TitleFontSizeMaxRule(BaseRule):
    rule_id = "title_font_max"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        for shape in slide.shapes:
            if not _is_title(shape) or not _has_text(shape):
                continue
            sizes = [s for s in _all_run_sizes(shape) if s > TITLE_MAX_PT]
            if sizes:
                pt = round(max(sizes), 1)
                return [Issue(
                    rule_id=self.rule_id,
                    severity="warning",
                    auto_fixable=False,
                    message=f"タイトルのフォントサイズが大きすぎます（最大: {pt}pt）",
                    suggestion=f"{TITLE_MAX_PT}pt 以下を推奨します",
                    details={"current": pt, "maximum": TITLE_MAX_PT},
                )]
        return []


class BodyFontSizeRule(BaseFixableRule):
    rule_id = "body_font_size"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        min_pt = None
        for shape in slide.shapes:
            if _is_title(shape) or not _has_text(shape):
                continue
            for pt in _all_run_sizes(shape):
                if pt < BODY_MIN_PT:
                    min_pt = pt if min_pt is None else min(min_pt, pt)
        if min_pt is not None:
            return [Issue(
                rule_id=self.rule_id,
                severity="warning",
                auto_fixable=True,
                message=f"本文のフォントサイズが小さすぎます（最小: {round(min_pt, 1)}pt）",
                suggestion=f"{BODY_MIN_PT}pt 以上を推奨します",
                details={"current": round(min_pt, 1), "minimum": BODY_MIN_PT},
            )]
        return []

    def fix(self, slide, presentation, style: Style) -> None:
        for shape in slide.shapes:
            if _is_title(shape) or not _has_text(shape):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None and run.font.size < Pt(BODY_MIN_PT):
                        run.font.size = Pt(BODY_MIN_PT)


class FontFamilyRule(BaseFixableRule):
    rule_id = "font_family"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        bad_fonts: set[str] = set()
        for shape in slide.shapes:
            if not _has_text(shape):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    name = run.font.name
                    if name and not any(name.startswith(p) for p in RECOMMENDED_PREFIXES):
                        bad_fonts.add(name)
        if bad_fonts:
            names = "、".join(sorted(bad_fonts))
            return [Issue(
                rule_id=self.rule_id,
                severity="warning",
                auto_fixable=True,
                message=f"推奨外のフォントが使用されています（{names}）",
                suggestion="メイリオ または BIZ UDPゴシック を推奨します",
                details={"fonts": sorted(bad_fonts)},
            )]
        return []

    def fix(self, slide, presentation, style: Style) -> None:
        for shape in slide.shapes:
            if not _has_text(shape):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and not any(run.font.name.startswith(p) for p in RECOMMENDED_PREFIXES):
                        run.font.name = FALLBACK_FONT


class FontVarietyRule(BaseRule):
    rule_id = "font_variety"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        fonts: set[str] = set()
        for shape in slide.shapes:
            if not _has_text(shape):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
        if len(fonts) > 3:
            return [Issue(
                rule_id=self.rule_id,
                severity="info",
                auto_fixable=False,
                message=f"フォント種類が多すぎます（{len(fonts)}種類）",
                suggestion="フォントは3種類以内に統一してください",
                details={"fonts": list(fonts), "count": len(fonts)},
            )]
        return []
