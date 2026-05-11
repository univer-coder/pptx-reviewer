from pptx.dml.color import RGBColor
from .base import BaseRule, BaseFixableRule, Issue, Style

BG_COLOR_STANDARD = RGBColor(0xFF, 0xFF, 0xFF)  # 白


def _slide_dimensions(presentation):
    return presentation.slide_width, presentation.slide_height


def _get_title_shape(slide):
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                return shape
    except Exception:
        pass
    return None


class TitlePositionRule(BaseFixableRule):
    rule_id = "title_position"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        _, height = _slide_dimensions(presentation)
        shape = _get_title_shape(slide)
        if shape is None:
            return []
        if shape.top > height * 0.25:
            return [Issue(
                rule_id=self.rule_id,
                severity="warning",
                auto_fixable=True,
                message="タイトルがスライド上部に配置されていません",
                suggestion="タイトルはスライド最上部（上から25%以内）に配置してください",
                details={"top_ratio": round(shape.top / height, 2)},
            )]
        return []

    def fix(self, slide, presentation, style: Style) -> None:
        _, height = _slide_dimensions(presentation)
        shape = _get_title_shape(slide)
        if shape is None:
            return
        if shape.top > height * 0.25:
            shape.top = int(height * 0.02)


class BackgroundColorRule(BaseFixableRule):
    rule_id = "background_color"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        if style != Style.STANDARD:
            return []
        try:
            fill = slide.background.fill
            rgb = fill.fore_color.rgb
            if (rgb.red, rgb.green, rgb.blue) != (0xFF, 0xFF, 0xFF):
                return [Issue(
                    rule_id=self.rule_id,
                    severity="info",
                    auto_fixable=True,
                    message=f"背景色が推奨色（#FFFFFF）ではありません",
                    suggestion="Standard スタイルでは背景色を白（#FFFFFF）に設定してください",
                    details={"current": f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"},
                )]
        except Exception:
            return [Issue(
                rule_id=self.rule_id,
                severity="info",
                auto_fixable=True,
                message="背景色が設定されていません",
                suggestion="Standard スタイルでは背景色を白（#FFFFFF）に設定してください",
                details={},
            )]
        return []

    def fix(self, slide, presentation, style: Style) -> None:
        if style != Style.STANDARD:
            return
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR_STANDARD


class ContainerMissingRule(BaseRule):
    rule_id = "container_missing"

    def check(self, slide, presentation, style: Style) -> list[Issue]:
        if style != Style.STANDARD:
            return []
        # 薄いグレー系の塗りつぶし図形（四角・角丸問わず）を検出
        has_container = False
        for shape in slide.shapes:
            try:
                fill = shape.fill
                rgb = fill.fore_color.rgb
                r, g, b = rgb.red, rgb.green, rgb.blue
                # グレー系（R=G=B かつ 180〜240 の薄いグレー）
                if abs(r - g) <= 10 and abs(g - b) <= 10 and 180 <= r <= 240:
                    has_container = True
                    break
            except Exception:
                continue

        if not has_container:
            return [Issue(
                rule_id=self.rule_id,
                severity="info",
                auto_fixable=False,
                message="グレーのコンテナブロックが配置されていません",
                suggestion="情報の塊を薄いグレー（#F2F2F2 など）の四角形で囲んでください",
                details={},
            )]
        return []
