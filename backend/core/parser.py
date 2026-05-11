import io
from pptx import Presentation
from .rules.base import Style, Issue
from .rules.font import BodyFontSizeRule, FontFamilyRule, FontVarietyRule
from .rules.density import TextLengthRule, LineCountRule
from .rules.contrast import ColorContrastRule, ThreeColorRule
from .rules.layout import TitlePositionRule, BackgroundColorRule, ContainerMissingRule
from .rules.structure import MissingTitleRule, TitleLinesRule, SlideCountRule, Chart3DRule, PieChartRule

SEVERITY_PENALTY = {"error": 10, "warning": 5, "info": 2}

ALL_RULES = [
    BodyFontSizeRule(),
    FontFamilyRule(),
    FontVarietyRule(),
    TextLengthRule(),
    LineCountRule(),
    ColorContrastRule(),
    ThreeColorRule(),
    TitlePositionRule(),
    BackgroundColorRule(),
    ContainerMissingRule(),
    MissingTitleRule(),
    TitleLinesRule(),
    SlideCountRule(),
    Chart3DRule(),
    PieChartRule(),
]

PRESENTATION_LEVEL_RULES = {"slide_count"}


def _calc_score(issues: list[Issue]) -> int:
    penalty = sum(SEVERITY_PENALTY[i.severity] for i in issues)
    return max(0, 100 - penalty)


def analyze(pptx_bytes: bytes, style: Style) -> dict:
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_results = []
    seen_presentation_rules: set[str] = set()

    for idx, slide in enumerate(prs.slides, start=1):
        issues: list[Issue] = []
        for rule in ALL_RULES:
            if rule.rule_id in PRESENTATION_LEVEL_RULES:
                if rule.rule_id in seen_presentation_rules:
                    continue
                seen_presentation_rules.add(rule.rule_id)
            try:
                issues.extend(rule.check(slide, prs, style))
            except Exception:
                pass

        slide_results.append({
            "slide_number": idx,
            "score": _calc_score(issues),
            "issues": [
                {
                    "rule_id": i.rule_id,
                    "severity": i.severity,
                    "auto_fixable": i.auto_fixable,
                    "message": i.message,
                    "suggestion": i.suggestion,
                    "details": i.details,
                }
                for i in issues
            ],
        })

    scores = [s["score"] for s in slide_results]
    overall_score = round(sum(scores) / len(scores)) if scores else 100
    all_issues = [i for s in slide_results for i in s["issues"]]
    has_autofixable = any(i["auto_fixable"] for i in all_issues)

    return {
        "total_slides": len(prs.slides),
        "overall_score": overall_score,
        "style": style.value,
        "has_autofixable": has_autofixable,
        "slides": slide_results,
    }
